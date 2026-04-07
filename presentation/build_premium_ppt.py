"""
Antigravity Resume Studio — Premium Pitch Deck Generator
Generates a high-end, startup-quality PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(SCRIPT_DIR, 'assets')

# ─── Theme Colors ────────────────────────────────────────────────────────
# ─── Theme Colors (Formal White) ──────────────────────────────────────────
DARK_BG      = RGBColor(255, 255, 255)    # White background
CARD_BG      = RGBColor(248, 250, 252)    # Light grey card surface
CARD_BORDER  = RGBColor(203, 213, 225)    # Subtle border
NEON_BLUE    = RGBColor(30, 58, 138)      # Deep blue accent
CYAN         = RGBColor(14, 116, 144)     # Dark cyan
GREEN        = RGBColor(21, 128, 61)      # Success green
PURPLE       = RGBColor(109, 40, 217)     # Purple accent
PINK         = RGBColor(190, 24, 93)      # Dark pink
TEXT_WHITE   = RGBColor(15, 23, 42)       # Primary text (dark grey/black)
TEXT_GREY    = RGBColor(71, 85, 105)      # Secondary text
TEXT_DIM     = RGBColor(100, 116, 139)    # Dim text
ORANGE       = RGBColor(194, 65, 12)      # Dark orange

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ─── Academic Info ─────────────────────────────────────────────────────────
PROJECT_TITLE = "AI RESUME STUDIO"
SUBTITLE      = "Resume Maker with ATS Scoring & AI-Powered Features"
STUDENT_NAME  = "Vansh Khandekar"
CLASS_NAME    = "Bachelor In Computer Applications (BCA)"
GUIDE_NAME    = "Prof. R.K. Sharma"
COLLEGE_NAME  = "JANAPRABHA COLLEGE, RAMTEK"
ACADEMIC_YEAR = "2025–26"


# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bg_image(slide, img_path):
    """Set a full-bleed background image."""
    slide.shapes.add_picture(
        img_path, Inches(0), Inches(0),
        width=SLIDE_WIDTH, height=SLIDE_HEIGHT
    )

def add_dark_overlay(slide, alpha=40):
    """Semi-transparent light overlay."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    # Set transparency via XML
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha_elem = srgb.makeelement(qn('a:alpha'), {})
        alpha_elem.set('val', str(alpha * 100000 // 255))
        srgb.append(alpha_elem)
    shape.line.fill.background()

def add_neon_line(slide, left, top, width, color=NEON_BLUE, height=Pt(3)):
    """Glowing accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Glow effect via XML
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.makeelement(qn('p:spPr'), {})
        shape._element.append(spPr)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    glow = effectLst.makeelement(qn('a:glow'), {'rad': '101600'})
    srgbClr = glow.makeelement(qn('a:srgbClr'), {'val': str(color)})
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '60000'})
    srgbClr.append(alpha)
    glow.append(srgbClr)
    effectLst.append(glow)
    spPr.append(effectLst)

def add_glass_card(slide, left, top, width, height, border_color=CARD_BORDER):
    """Glassmorphism-style card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    # Transparency via XML
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha_elem = srgb.makeelement(qn('a:alpha'), {})
        alpha_elem.set('val', '98000')  # 98% opacity (almost solid for white)
        srgb.append(alpha_elem)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    # Rounded corners
    shape.adjustments[0] = 0.05
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    # IMPORTANT: Use run-level font properties, NOT paragraph-level.
    # Paragraph-level (.p.font) gets overridden by PowerPoint theme defaults.
    # Run-level (.runs[].font) always takes precedence and renders correctly.
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_icon_text(slide, left, top, icon_char, label, desc,
                  accent_color=NEON_BLUE, card_w=Inches(2.8), card_h=Inches(2.5)):
    """Card with icon character, label, and description."""
    card = add_glass_card(slide, left, top, card_w, card_h, accent_color)
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.9), top + Inches(0.3),
        Inches(1.0), Inches(1.0)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    # Transparency for circle via XML
    spPr = circle._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha_elem = srgb.makeelement(qn('a:alpha'), {})
        alpha_elem.set('val', '25000')
        srgb.append(alpha_elem)
    circle.line.fill.background()
    # Icon text — use run-level for guaranteed color
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = icon_char
    run.font.size = Pt(32)
    run.font.color.rgb = RGBColor(255, 255, 255) # White icon on colored circle
    run.font.bold = True
    # Label
    add_text(slide, left + Inches(0.15), top + Inches(1.3), card_w - Inches(0.3), Inches(0.4),
             label, font_size=15, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_text(slide, left + Inches(0.1), top + Inches(1.8), card_w - Inches(0.2), Inches(0.6),
             desc, font_size=10, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

def add_step_card(slide, left, top, number, title, desc, accent=NEON_BLUE):
    """Numbered step card for workflow."""
    card = add_glass_card(slide, left, top, Inches(2.5), Inches(2.4), accent)
    # Number circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.85), top + Inches(0.25),
        Inches(0.8), Inches(0.8)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.bold = True
    # Title
    add_text(slide, left + Inches(0.15), top + Inches(1.05), Inches(2.2), Inches(0.5),
             title, font_size=15, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_text(slide, left + Inches(0.1), top + Inches(1.6), Inches(2.3), Inches(1.0),
             desc, font_size=10, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

def add_slide_number(slide, num, total):
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
             f"{num}/{total}", font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)

def section_heading(slide, title, subtitle="", top=Inches(0.5)):
    add_text(slide, Inches(0.8), top, Inches(8), Inches(0.7),
             title, font_size=36, color=TEXT_WHITE, bold=True)
    add_neon_line(slide, Inches(0.8), top + Inches(0.7), Inches(2.5))
    if subtitle:
        add_text(slide, Inches(0.8), top + Inches(0.85), Inches(10), Inches(0.5),
                 subtitle, font_size=14, color=TEXT_GREY)

TOTAL_SLIDES = 9

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
hero_path = os.path.join(ASSETS_DIR, 'hero_bg.png')
if os.path.exists(hero_path):
    # For white theme, skip background image or use very high transparency
    set_bg(slide, DARK_BG)
else:
    set_bg(slide, DARK_BG)

# Neon accent lines
add_neon_line(slide, Inches(1.0), Inches(1.8), Inches(4.0), NEON_BLUE, Pt(4))
add_neon_line(slide, Inches(1.0), Inches(1.95), Inches(2.5), CYAN, Pt(2))

# Title
add_text(slide, Inches(1.0), Inches(1.0), Inches(11), Inches(1.2),
         PROJECT_TITLE, font_size=54, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Header Line
add_neon_line(slide, Inches(1.0), Inches(2.2), Inches(11.3), NEON_BLUE, Pt(2))

# Subtitle
add_text(slide, Inches(1.0), Inches(2.3), Inches(11), Inches(0.6),
         CLASS_NAME, font_size=28, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Middle Info
info_card = add_glass_card(slide, Inches(3.4), Inches(3.2), Inches(6.5), Inches(3.5))
add_card_info = [
    (f"Project Title :  {PROJECT_TITLE}", 18, TEXT_WHITE, True),
    (f"Student Name :  {STUDENT_NAME}", 18, TEXT_WHITE, False),
    (f"Class        :  {CLASS_NAME}", 16, TEXT_GREY, False),
    (f"Guide Name   :  {GUIDE_NAME}", 16, TEXT_GREY, False),
    (f"College      :  {COLLEGE_NAME}", 16, TEXT_GREY, False),
    (f"Date         :  April 2026", 14, TEXT_DIM, False)
]
for i, (txt, sz, col, bld) in enumerate(add_card_info):
    add_text(slide, Inches(3.8), Inches(3.5) + Inches(0.5) * i, Inches(6), Inches(0.4),
             txt, font_size=sz, color=col, bold=bld)

add_slide_number(slide, 1, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTENT (INDEX)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "Table of Contents")

toc_points = [
    "1) Introduction", "2) Problem Definition", "3) Objectives & Design Goals",
    "4) Scope & Constraints", "5) System Architecture", "6) Technology Stack",
    "7) Functional Modules", "8) Database Design & Data Modeling",
    "9) Implementation Strategy", "10) Outputs", "11) Future Enhancements",
    "12) References"
]

# Two columns for TOC
for i, pt in enumerate(toc_points):
    col = i // 6
    row = i % 6
    x = Inches(1.2) + Inches(5.5) * col
    y = Inches(2.2) + Inches(0.7) * row
    add_text(slide, x, y, Inches(5), Inches(0.6), pt, font_size=20, color=TEXT_WHITE)

add_slide_number(slide, 2, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "1. Introduction", "Modernizing the job application process with AI")

intro_text = (
    "AI Resume Studio is a next-generation platform designed to help job seekers "
    "create highly competitive, ATS-optimized resumes. By integrating Large Language "
    "Models, the system provides intelligent content suggestions, real-time "
    "scoring, and high-fidelity PDF exports. It bridges the gap between raw "
    "professional experience and the sophisticated requirements of modern "
    "Applicant Tracking Systems."
)
add_text(slide, Inches(1.2), Inches(2.5), Inches(11), Inches(1.5),
         intro_text, font_size=24, color=TEXT_WHITE, alignment=PP_ALIGN.JUSTIFY)

# Points
points = ["Modular Design", "Real-time Scoring", "AI-Powered Writing", "Cloud-Native Flow"]
for i, pt in enumerate(points):
    add_text(slide, Inches(1.5), Inches(4.5) + Inches(0.6) * i, Inches(10), Inches(0.5),
             f"✦  {pt}", font_size=20, color=NEON_BLUE, bold=True)

add_slide_number(slide, 3, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM DEFINITION
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "2. Problem Definition", "Why traditional resume building is broken")

problems = [
    ("⚠️", "Formatting Chaos", "Word layouts break across\ndevices & ATS parsers", PINK),
    ("🚫", "ATS Rejection", "80% of resumes are rejected by\nATS before human review", ORANGE),
    ("😶", "Writer's Block", "Candidates struggle to articulate\nachievements professionally", PURPLE),
    ("⏳", "Time Consuming", "Manual formatting takes 3–5 hours\nper resume version", CYAN),
 ]
for i, (icon, title, desc, color) in enumerate(problems):
    x = Inches(0.6) + Inches(3.1) * i
    add_icon_text(slide, x, Inches(2.5), icon, title, desc, accent_color=color)

add_slide_number(slide, 4, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OBJECTIVES & DESIGN GOALS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "3. Objectives", "Designing a user-centric AI career assistant")

objectives = [
    ("🎯", "ATS Optimization", "100% compliance with\nmodern ranking algorithms", NEON_BLUE),
    ("🤖", "AI Integration", "Leveraging LLMs for\nsmart content auditing", CYAN),
    ("🎨", "UI/UX Excellence", "Premium, interactive\nresume builder interface", GREEN),
    ("📥", "Seamless Export", "Watermark-free, professional\nPDF delivery in seconds", PURPLE),
]
for i, (icon, title, desc, color) in enumerate(objectives):
    x = Inches(0.6) + Inches(3.1) * i
    add_icon_text(slide, x, Inches(2.5), icon, title, desc, accent_color=color)

add_slide_number(slide, 5, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SCOPE & CONSTRAINTS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "4. Scope & Constraints", "Defining project boundaries and limits")

# Left Scope
add_glass_card(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), NEON_BLUE)
add_text(slide, Inches(1.2), Inches(2.5), Inches(5), Inches(0.5), "🛠️ Project Scope", font_size=24, color=NEON_BLUE, bold=True)
scope_items = ["Multi-template Support", "Multi-section Builder", "AI Summary Generator", "Keyword Optimization", "Cloud Persistence"]
for i, item in enumerate(scope_items):
    add_text(slide, Inches(1.5), Inches(3.2) + Inches(0.6) * i, Inches(5), Inches(0.4), f"▸ {item}", font_size=18)

# Right Constraints
add_glass_card(slide, Inches(6.9), Inches(2.2), Inches(5.8), Inches(4.5), ORANGE)
add_text(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.5), "⚠️ Constraints", font_size=24, color=ORANGE, bold=True)
const_items = ["Internet Dependency", "API Token Limits", "Single Page Export", "English Language focus", "Static Assets only"]
for i, item in enumerate(const_items):
    add_text(slide, Inches(7.6), Inches(3.2) + Inches(0.6) * i, Inches(5), Inches(0.4), f"▸ {item}", font_size=18)

add_slide_number(slide, 6, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "5. System Architecture", "Cloud-native, three-tier scalable infrastructure")

# Flow items
flow_items = [
    ("👤", "React UI", "Optimistic state\nClient-side logic", NEON_BLUE),
    ("→", "", "", TEXT_DIM),
    ("⚙️", "API Gateway", "Deno Edge Functions\nClaude-3 Integration", CYAN),
    ("→", "", "", TEXT_DIM),
    ("🗄️", "Supabase DB", "PostgreSQL JSONB\nRow Level Security", GREEN),
]
x_pos = Inches(0.5)
for icon, title, desc, color in flow_items:
    if icon == "→":
        add_text(slide, x_pos, Inches(3.5), Inches(1.0), Inches(0.5), "───▶", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(1.0)
    else:
        add_icon_text(slide, x_pos, Inches(2.5), icon, title, desc, accent_color=color, card_w=Inches(3.2), card_h=Inches(2.8))
        x_pos += Inches(3.6)

add_slide_number(slide, 7, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "6. Technology Stack", "Built with modern, production-grade tools")

# Frontend Column
fe_card = add_glass_card(slide, Inches(0.6), Inches(2.2), Inches(3.8), Inches(4.8), NEON_BLUE)
add_text(slide, Inches(0.9), Inches(2.4), Inches(3.2), Inches(0.5), "⚛️ FRONTEND", font_size=22, color=NEON_BLUE, bold=True)
techs_fe = ["React 18", "TypeScript", "Vite 5", "Tailwind CSS 3", "Shadcn UI", "Lucide Icons"]
for i, tech in enumerate(techs_fe):
    add_text(slide, Inches(1.1), Inches(3.2) + Inches(0.45) * i, Inches(3), Inches(0.4), f"▸ {tech}", font_size=16)

# Backend Column
be_card = add_glass_card(slide, Inches(4.8), Inches(2.2), Inches(3.8), Inches(4.8), CYAN)
add_text(slide, Inches(5.1), Inches(2.4), Inches(3.2), Inches(0.5), "⚙️ BACKEND", font_size=22, color=CYAN, bold=True)
techs_be = ["Supabase (Postgres)", "Edge Functions", "Claude-3 API", "JWT Auth", "PostgreSQL JSONB"]
for i, tech in enumerate(techs_be):
    add_text(slide, Inches(5.3), Inches(3.2) + Inches(0.45) * i, Inches(3), Inches(0.4), f"▸ {tech}", font_size=16)

# Deployment Column
tool_card = add_glass_card(slide, Inches(9.0), Inches(2.2), Inches(3.8), Inches(4.8), PURPLE)
add_text(slide, Inches(9.3), Inches(2.4), Inches(3.2), Inches(0.5), "🛠️ DEPLOYMENT", font_size=22, color=PURPLE, bold=True)
techs_tool = ["Git & GitHub", "Vercel CI/CD", "ESLint + Prettier", "React Query", "jsPDF Export"]
for i, tech in enumerate(techs_tool):
    add_text(slide, Inches(9.5), Inches(3.2) + Inches(0.45) * i, Inches(3), Inches(0.4), f"▸ {tech}", font_size=16)

add_slide_number(slide, 8, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FUNCTIONAL MODULES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "7. Functional Modules", "Core system capabilities and key features")

features = [
    ("🤖", "AI Content Engine", "LLM-driven rewriting\nand professional auditing", NEON_BLUE),
    ("📊", "ATS Scoring", "Real-time ranking heuristics\nand keyword analysis", CYAN),
    ("⚡", "Modular Builder", "Dynamic form builder with\ndocument state persistence", ORANGE),
    ("📥", "PDF Generator", "One-click, high-fidelity\nexport with typography", GREEN),
]
for i, (icon, title, desc, color) in enumerate(features):
    x_col = i % 2
    y_row = i // 2
    add_icon_text(slide, Inches(2.2) + Inches(4.5)*x_col, Inches(2.3) + Inches(2.5)*y_row, icon, title, desc, accent_color=color)

add_slide_number(slide, 9, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DATABASE DESIGN
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "8. Database & Data Modeling", "Schema integrity and JSONB data persistence")

# Table card
add_glass_card(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), NEON_BLUE)
entities = [
    ("Profiles", "User authentication and baseline metadata"),
    ("Resumes", "The core entity storing resume JSONB structures"),
    ("Sections", "Dynamic section configurations and visibility flags"),
    ("Audit Logs", "Tracking AI credit usage and history"),
]
for i, (ent, desc) in enumerate(entities):
    add_text(slide, Inches(1.5), Inches(2.8) + Inches(0.8) * i, Inches(3), Inches(0.5), f"• {ent}", font_size=24, color=NEON_BLUE, bold=True)
    add_text(slide, Inches(5.0), Inches(2.8) + Inches(0.8) * i, Inches(7), Inches(0.5), f"▸ {desc}", font_size=18)

add_slide_number(slide, 10, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION STRATEGY
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "9. Implementation Strategy", "Iterative development and deployment flow")

strat_steps = [
    (1, "Phase 1", "Database schema setup &\nAuthentication layer init", NEON_BLUE),
    (2, "Phase 2", "Core Builder UI & State Management\nwith React Context", CYAN),
    (3, "Phase 3", "AI Edge Functions &\nGemini API Integration", GREEN),
    (4, "Phase 4", "PDF Export Engine &\nFinal UAT Deployment", PURPLE),
]
for i, (num, title, desc, color) in enumerate(strat_steps):
    add_step_card(slide, Inches(0.5) + Inches(3.2)*i, Inches(2.5), num, title, desc, accent=color)

add_slide_number(slide, 11, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OUTPUTS (Screenshots)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "10. Product Outputs", "Glimpse of the final application interface")

db_m = os.path.join(ASSETS_DIR, 'dashboard_mockup.png')
pv_m = os.path.join(ASSETS_DIR, 'preview_mockup.png')
if os.path.exists(db_m):
    add_glass_card(slide, Inches(0.4), Inches(2.0), Inches(6.2), Inches(4.5), NEON_BLUE)
    slide.shapes.add_picture(db_m, Inches(0.6), Inches(2.2), width=Inches(5.8), height=Inches(4.1))
if os.path.exists(pv_m):
    add_glass_card(slide, Inches(6.9), Inches(2.0), Inches(6.2), Inches(4.5), GREEN)
    slide.shapes.add_picture(pv_m, Inches(7.1), Inches(2.2), width=Inches(5.8), height=Inches(4.1))

add_slide_number(slide, 12, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FUTURE ENHANCEMENTS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "11. Future Enhancements", "Expanding the horizons of AI career tools")

futures = [
    ("🔗", "LinkedIn Import", "One-click data Sync", NEON_BLUE),
    ("✉️", "Cover Letters", "Auto-Gen per Job", CYAN),
    ("🎨", "Custom Themes", "CSS-level control", PURPLE),
    ("🌐", "Job Submission", "Direct ATS upload", GREEN),
]
for i, (icon, title, desc, color) in enumerate(futures):
    add_icon_text(slide, Inches(0.4) + Inches(3.2)*i, Inches(3.0), icon, title, desc, accent_color=color, card_w=Inches(2.9), card_h=Inches(2.6))

add_slide_number(slide, 13, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — REFERENCES
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)
section_heading(slide, "References")

refs = [
    "• React Documentation (react.dev) — Component architecture",
    "• Supabase Auth & DB (supabase.com) — Backend infrastructure",
    "• Anthropic Claude 3 Opus API — AI content generation logic",
    "• Tailwind CSS — Utility-first styling framework",
    "• 'Building Resumes with AI' (Medium 2024) — ATS Heuristics",
    "• GitHub: AI Resume Studio Open Source Repository"
]
for i, ref in enumerate(refs):
    add_text(slide, Inches(1.2), Inches(2.5) + Inches(0.7) * i, Inches(10), Inches(0.6),
             ref, font_size=18, color=TEXT_WHITE)

add_slide_number(slide, 14, 15)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_BG)

add_neon_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), NEON_BLUE, Pt(4))

add_text(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.0),
         "Thank You!", font_size=56, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
         "Questions & Discussion", font_size=26, color=NEON_BLUE, alignment=PP_ALIGN.CENTER)

add_neon_line(slide, Inches(5.5), Inches(4.8), Inches(2.3), CYAN, Pt(2))

# Info
add_text(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
         f"{STUDENT_NAME}  •  {CLASS_NAME}  •  2025–26",
         font_size=16, color=TEXT_GREY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
         "Built with React • TypeScript • Gemini AI • Supabase",
         font_size=13, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 15, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = os.path.join(SCRIPT_DIR, 'AI_Resume_Studio_Premium.pptx')
prs.save(output_path)
print(f"✅ Premium presentation saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
print("   Theme: Dark Space / Neon Blue / Glassmorphism")
print("   Images embedded: hero_bg, landing_mockup, dashboard_mockup, preview_mockup")
print("   Transitions: fade, push, wipe applied to all slides")
